import os
import io
import email
from pathlib import Path
import pandas as pd
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
import PyPDF2
from tika import parser
from utils.logging_config import logger
from config.settings import settings
from typing import Dict, Any, Optional

class DocumentLoader:
    def __init__(self):
        self.supported = {
            '.pdf': self._load_pdf,
            '.docx': self._load_docx,
            '.pptx': self._load_pptx,
            '.csv': self._load_csv,
            '.txt': self._load_text,
            '.png': self._load_image,
            '.jpg': self._load_image,
            '.jpeg': self._load_image,
            '.eml': self._load_email
        }

    def load_document(self, file_path: str, file_content: Optional[bytes] = None) -> Dict[str, Any]:
        ext = Path(file_path).suffix.lower()
        if ext not in self.supported:
            raise ValueError(f"Unsupported: {ext}")
        loader = self.supported[ext]
        logger.info(f"Loading {file_path}")
        content = loader(file_path=file_path) if file_content is None else loader(file_content=file_content)
        return {
            "filename": Path(file_path).name,
            "file_type": ext,
            "content": content,
            "file_size": len(file_content) if file_content else os.path.getsize(file_path),
            "metadata": {"source": file_path, "type": ext, "loader": loader.__name__}
        }

    def _load_pdf(self, file_path=None, file_content=None) -> str:
        try:
            if file_content is not None:
                try:
                    pdf = PyPDF2.PdfReader(io.BytesIO(file_content))
                    return "\n".join(p.extract_text() or "" for p in pdf.pages)
                except Exception:
                    result = parser.from_buffer(file_content)
                    if isinstance(result, dict):
                        return result.get("content", "") or ""
                    elif isinstance(result, tuple) and isinstance(result[1], dict):
                        return result[1].get("content", "") or ""
                    return ""
            elif file_path is not None:
                try:
                    pdf = PyPDF2.PdfReader(file_path)
                    return "\n".join(p.extract_text() or "" for p in pdf.pages)
                except Exception:
                    result = parser.from_file(file_path)
                    if isinstance(result, dict):
                        return result.get("content", "") or ""
                    elif isinstance(result, tuple) and isinstance(result[1], dict):
                        return result[1].get("content", "") or ""
                    return ""
            else:
                return ""
        except Exception as e:
            logger.error(f"PDF load error: {e}")
            return ""

    def _load_docx(self, file_path=None, file_content=None) -> str:
        try:
            doc = Document(io.BytesIO(file_content)) if file_content is not None else Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs)
            for tbl in doc.tables:
                for row in tbl.rows:
                    text += "\n" + "\t".join(cell.text for cell in row.cells)
            return text
        except Exception as e:
            logger.error(f"DOCX load error: {e}")
            return ""

    def _load_pptx(self, file_path=None, file_content=None) -> str:
        try:
            prs = Presentation(io.BytesIO(file_content)) if file_content is not None else Presentation(file_path)
            return "\n".join(
                getattr(shape, "text", "")
                for s in prs.slides for shape in s.shapes
                if hasattr(shape, "text") and isinstance(getattr(shape, "text", None), str)
            )
        except Exception as e:
            logger.error(f"PPTX load error: {e}")
            return ""

    def _load_csv(self, file_path=None, file_content=None) -> str:
        try:
            if file_content is not None:
                df = pd.read_csv(io.BytesIO(file_content))
            elif file_path is not None:
                df = pd.read_csv(file_path)
            else:
                return ""
            return df.to_string(index=False)
        except Exception as e:
            logger.error(f"CSV load error: {e}")
            return ""

    def _load_text(self, file_path=None, file_content=None) -> str:
        try:
            if file_content is not None:
                return file_content.decode('utf-8', errors='ignore')
            elif file_path is not None:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            else:
                return ""
        except Exception as e:
            logger.error(f"TXT load error: {e}")
            return ""

    def _load_image(self, file_path=None, file_content=None) -> str:
        try:
            if file_content is not None:
                img = Image.open(io.BytesIO(file_content))
            elif file_path is not None:
                img = Image.open(file_path)
            else:
                return ""
            return pytesseract.image_to_string(img)
        except Exception as e:
            logger.error(f"Image load error: {e}")
            return ""

    def _load_email(self, file_path=None, file_content=None) -> str:
        try:
            if file_content is not None:
                msg = email.message_from_bytes(file_content)
            elif file_path is not None:
                with open(file_path, 'rb') as f:
                    msg = email.message_from_binary_file(f)
            else:
                return ""
            text = f"Subject: {msg.get('Subject', '')}\nFrom: {msg.get('From', '')}\nTo: {msg.get('To', '')}\nDate: {msg.get('Date', '')}\n\n"
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if isinstance(payload, bytes):
                            text += payload.decode('utf-8', errors='ignore')
                        elif isinstance(payload, str):
                            text += payload
            else:
                payload = msg.get_payload(decode=True)
                if isinstance(payload, bytes):
                    text += payload.decode('utf-8', errors='ignore')
                elif isinstance(payload, str):
                    text += payload
            return text
        except Exception as e:
            logger.error(f"EML load error: {e}")
            return ""

    def validate_file(self, file_path: str, file_size: int) -> bool:
        ext = Path(file_path).suffix.lower()
        return ext in self.supported and file_size <= settings.max_file_size

document_loader = DocumentLoader()