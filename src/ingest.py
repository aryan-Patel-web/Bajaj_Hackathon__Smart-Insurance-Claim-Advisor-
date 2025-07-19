import os
import fitz  # PyMuPDF
import docx2txt
import csv
import pptx
import logging
import base64
import pytesseract
import mimetypes
import traceback
from PIL import Image
from io import BytesIO
from typing import List, Dict, Optional
from langchain.docstore.document import Document
from langchain_community.document_loaders import PyPDFLoader, UnstructuredEmailLoader
from utils.logging_config import setup_logging

logger = logging.getLogger(__name__)
setup_logging()

class DocumentIngestor:
    def __init__(self):
        self.supported_extensions = [".pdf", ".docx", ".pptx", ".csv", ".eml", ".jpg", ".jpeg", ".png"]

    def read_file(self, file_path: str) -> Optional[str]:
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".pdf":
                return self.read_pdf(file_path)
            elif ext == ".docx":
                return self.read_docx(file_path)
            elif ext == ".pptx":
                return self.read_pptx(file_path)
            elif ext == ".csv":
                return self.read_csv(file_path)
            elif ext == ".eml":
                return self.read_email(file_path)
            elif ext in [".jpg", ".jpeg", ".png"]:
                return self.read_image(file_path)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return None
        except Exception as e:
            logger.error(f"Error reading {file_path}: {str(e)}")
            traceback.print_exc()
            return None

    def read_pdf(self, file_path: str) -> str:
        loader = PyPDFLoader(file_path)
        pages = loader.load_and_split()
        return "\n".join([p.page_content for p in pages])

    def read_docx(self, file_path: str) -> str:
        return docx2txt.process(file_path)

    def read_pptx(self, file_path: str) -> str:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def read_csv(self, file_path: str) -> str:
        with open(file_path, newline='', encoding='utf-8') as csvfile:
            return "\n".join([", ".join(row) for row in csv.reader(csvfile)])

    def read_image(self, file_path: str) -> str:
        img = Image.open(file_path)
        return pytesseract.image_to_string(img)

    def read_email(self, file_path: str) -> str:
        loader = UnstructuredEmailLoader(file_path)
        docs = loader.load()
        return "\n".join([doc.page_content for doc in docs])

    def process_upload(self, uploaded_file) -> List[Document]:
        try:
            extension = os.path.splitext(uploaded_file.name)[1].lower()
            mime_type, _ = mimetypes.guess_type(uploaded_file.name)

            if extension not in self.supported_extensions:
                raise ValueError("Unsupported file format")

            content_bytes = uploaded_file.read()
            file_path = f"/tmp/{uploaded_file.name}"

            with open(file_path, "wb") as f:
                f.write(content_bytes)

            text = self.read_file(file_path)
            os.remove(file_path)

            if not text:
                raise ValueError("Empty content or unsupported format")

            metadata = {
                "source": uploaded_file.name,
                "type": mime_type or extension
            }
            return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            traceback.print_exc()
            return []

    def base64_to_docs(self, base64_string: str, filename: str) -> List[Document]:
        try:
            extension = os.path.splitext(filename)[1].lower()
            mime_type, _ = mimetypes.guess_type(filename)

            if extension not in self.supported_extensions:
                raise ValueError("Unsupported file format")

            file_path = f"/tmp/{filename}"
            with open(file_path, "wb") as f:
                f.write(base64.b64decode(base64_string))

            text = self.read_file(file_path)
            os.remove(file_path)

            if not text:
                raise ValueError("Empty content or unsupported format")

            metadata = {
                "source": filename,
                "type": mime_type or extension
            }
            return [Document(page_content=text, metadata=metadata)]
        except Exception as e:
            logger.error(f"Error decoding base64 document: {str(e)}")
            traceback.print_exc()
            return []
