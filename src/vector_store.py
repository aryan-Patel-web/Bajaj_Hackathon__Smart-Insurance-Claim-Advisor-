import os
import logging
from typing import List, Dict, Any, Optional
import pandas as pd
from pathlib import Path
import tempfile
import hashlib
from dataclasses import dataclass
from sentence_transformers import SentenceTransformer
from cassandra.cluster import Cluster
from cassandra.auth import PlainTextAuthProvider
import uuid
from datetime import datetime
import json
import time
import requests
from urllib.parse import urljoin

# Document processing libraries
import PyPDF2
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import email
from email.mime.text import MIMEText
import zipfile

logger = logging.getLogger(__name__)

@dataclass
class DocumentChunk:
    """Represents a processed document chunk"""
    chunk_id: str
    content: str
    metadata: Dict[str, Any]
    embedding: Optional[List[float]] = None

class DocumentProcessor:
    """Handles multi-format document processing and ingestion"""

    def __init__(self, embedding_model_name: str = "sentence-transformers/all-MiniLM-L6-v2"):
        self.embedding_model = SentenceTransformer(embedding_model_name)
        self.chunk_size = 1000  # ~1KB chunks
        self.overlap = 200

    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF files"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""

    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX files"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""

    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX files"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and isinstance(getattr(shape, "text", None), str):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""

    def extract_text_from_csv(self, file_path: str) -> str:
        """Extract text from CSV files"""
        try:
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            return text
        except Exception as e:
            logger.error(f"Error extracting CSV text: {e}")
            return ""

    def extract_text_from_image(self, file_path: str) -> str:
        """Extract text from images using OCR"""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.error(f"Error extracting image text: {e}")
            return ""

    def extract_text_from_email(self, file_path: str) -> str:
        """Extract text from email files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                msg = email.message_from_file(file)
                text = ""
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            payload = part.get_payload(decode=True)
                            if isinstance(payload, bytes):
                                text += payload.decode('utf-8')
                            elif isinstance(payload, str):
                                text += payload
                else:
                    payload = msg.get_payload(decode=True)
                    if isinstance(payload, bytes):
                        text = payload.decode('utf-8')
                    elif isinstance(payload, str):
                        text = payload
                return text
        except Exception as e:
            logger.error(f"Error extracting email text: {e}")
            return ""

    def extract_text_from_file(self, file_path: str) -> str:
        """Route to appropriate extraction method based on file extension"""
        file_extension = Path(file_path).suffix.lower()

        extraction_methods = {
            '.pdf': self.extract_text_from_pdf,
            '.docx': self.extract_text_from_docx,
            '.pptx': self.extract_text_from_pptx,
            '.csv': self.extract_text_from_csv,
            '.jpg': self.extract_text_from_image,
            '.jpeg': self.extract_text_from_image,
            '.png': self.extract_text_from_image,
            '.eml': self.extract_text_from_email,
            '.txt': lambda path: open(path, 'r', encoding='utf-8').read()
        }

        if file_extension in extraction_methods:
            return extraction_methods[file_extension](file_path)
        else:
            logger.warning(f"Unsupported file type: {file_extension}")
            return ""

    def chunk_text(self, text: str, metadata: Dict[str, Any]) -> List[DocumentChunk]:
        """Split text into chunks with metadata"""
        chunks = []
        words = text.split()

        for i in range(0, len(words), self.chunk_size - self.overlap):
            chunk_words = words[i:i + self.chunk_size]
            chunk_text = ' '.join(chunk_words)

            # Create unique chunk ID
            chunk_id = hashlib.md5(f"{metadata['filename']}_{i}_{chunk_text[:50]}".encode()).hexdigest()

            chunk_metadata = {
                **metadata,
                'chunk_index': i // (self.chunk_size - self.overlap),
                'chunk_start': i,
                'chunk_end': min(i + self.chunk_size, len(words)),
                'timestamp': datetime.now().isoformat()
            }

            chunks.append(DocumentChunk(
                chunk_id=chunk_id,
                content=chunk_text,
                metadata=chunk_metadata
            ))

        return chunks

    def generate_embeddings(self, chunks: List[DocumentChunk]) -> List[DocumentChunk]:
        """Generate embeddings for document chunks"""
        texts = [chunk.content for chunk in chunks]
        embeddings = self.embedding_model.encode(texts)

        for chunk, embedding in zip(chunks, embeddings):
            chunk.embedding = embedding.tolist()

        return chunks

    def process_document(self, file_path: str, document_type: str = "insurance_policy") -> List[DocumentChunk]:
        """Process a single document and return chunks with embeddings"""
        try:
            # Extract text
            text = self.extract_text_from_file(file_path)
            if not text.strip():
                logger.warning(f"No text extracted from {file_path}")
                return []

            # Create metadata
            metadata = {
                'filename': Path(file_path).name,
                'filepath': file_path,
                'document_type': document_type,
                'file_size': os.path.getsize(file_path),
                'processed_at': datetime.now().isoformat(),
                'text_length': len(text)
            }

            # Chunk text
            chunks = self.chunk_text(text, metadata)

            # Generate embeddings
            chunks = self.generate_embeddings(chunks)

            logger.info(f"Processed {len(chunks)} chunks from {Path(file_path).name}")
            return chunks

        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
            return []

class AstraDBIngester:
    """Handles ingestion of processed chunks into Astra DB"""

    def __init__(self, token: str, database_id: str, region: str = "us-east1", keyspace: str = None):
        self.token = token
        self.database_id = database_id
        self.region = region
        self.keyspace = keyspace or "default_keyspace"
        self.session = None
        self.cluster = None
        self.connect()

    def get_secure_connect_bundle(self) -> str:
        """Download secure connect bundle if not present"""
        bundle_path = f"secure-connect-{self.database_id}.zip"
        
        if not os.path.exists(bundle_path):
            try:
                # Download the secure connect bundle
                url = f"https://datastax-cluster-config-prod.s3.us-west-2.amazonaws.com/{self.database_id}/secure-connect-{self.database_id}.zip"
                
                # Alternative: Use Astra DB REST API to get the bundle
                headers = {
                    'Authorization': f'Bearer {self.token}',
                    'Content-Type': 'application/json'
                }
                
                # First, try to get database info
                db_info_url = f"https://api.astra.datastax.com/v2/databases/{self.database_id}"
                response = requests.get(db_info_url, headers=headers)
                
                if response.status_code == 200:
                    db_info = response.json()
                    # Get the secure bundle download URL
                    bundle_url = f"https://api.astra.datastax.com/v2/databases/{self.database_id}/secureBundleURL"
                    bundle_response = requests.post(bundle_url, headers=headers)
                    
                    if bundle_response.status_code == 200:
                        download_url = bundle_response.json()['downloadURL']
                        
                        # Download the bundle
                        bundle_data = requests.get(download_url)
                        with open(bundle_path, 'wb') as f:
                            f.write(bundle_data.content)
                        
                        logger.info(f"Downloaded secure connect bundle: {bundle_path}")
                    else:
                        logger.error(f"Failed to get bundle URL: {bundle_response.text}")
                        return None
                else:
                    logger.error(f"Failed to get database info: {response.text}")
                    return None
                    
            except Exception as e:
                logger.error(f"Error downloading secure connect bundle: {e}")
                return None
        
        return bundle_path

    def connect(self):
        """Establish connection to Astra DB"""
        try:
            # Get or download secure connect bundle
            bundle_path = self.get_secure_connect_bundle()
            if not bundle_path:
                logger.error("Failed to get secure connect bundle")
                return

            cloud_config = {
                'secure_connect_bundle': bundle_path
            }

            auth_provider = PlainTextAuthProvider('token', self.token)
            self.cluster = Cluster(cloud=cloud_config, auth_provider=auth_provider)
            self.session = self.cluster.connect()

            if self.session:
                # List available keyspaces
                keyspaces_result = self.session.execute("SELECT keyspace_name FROM system_schema.keyspaces")
                available_keyspaces = [row.keyspace_name for row in keyspaces_result]
                logger.info(f"Available keyspaces: {available_keyspaces}")

                # Use existing keyspace or create new one
                if self.keyspace not in available_keyspaces:
                    logger.info(f"Creating keyspace: {self.keyspace}")
                    self.session.execute(f"""
                        CREATE KEYSPACE IF NOT EXISTS {self.keyspace}
                        WITH replication = {{'class': 'SimpleStrategy', 'replication_factor': 1}}
                    """)
                    # Wait for keyspace creation
                    time.sleep(5)

                self.session.set_keyspace(self.keyspace)

                # Create tables
                self.create_tables()

                logger.info(f"Connected to Astra DB successfully using keyspace: {self.keyspace}")
            else:
                logger.error("Failed to connect to Astra DB: session is None")
        except Exception as e:
            logger.error(f"Failed to connect to Astra DB: {e}")
            raise

    def create_tables(self):
        """Create necessary tables for document storage"""
        try:
            if self.session:
                # Document chunks table
                self.session.execute("""
                    CREATE TABLE IF NOT EXISTS document_chunks (
                        chunk_id text PRIMARY KEY,
                        content text,
                        embedding list<float>,
                        metadata text,
                        document_type text,
                        filename text,
                        created_at timestamp
                    )
                """)

                # Conversation history table
                self.session.execute("""
                    CREATE TABLE IF NOT EXISTS conversation_history (
                        conversation_id text,
                        message_id text,
                        timestamp timestamp,
                        user_message text,
                        bot_response text,
                        context text,
                        PRIMARY KEY (conversation_id, message_id)
                    )
                """)

                # Document metadata table
                self.session.execute("""
                    CREATE TABLE IF NOT EXISTS document_metadata (
                        document_id text PRIMARY KEY,
                        filename text,
                        document_type text,
                        file_size bigint,
                        processed_at timestamp,
                        total_chunks int,
                        processing_status text
                    )
                """)

                logger.info("Tables created successfully")
            else:
                logger.error("Cannot create tables: session is None")
        except Exception as e:
            logger.error(f"Error creating tables: {e}")
            raise

    def ingest_chunks(self, chunks: List[DocumentChunk]) -> bool:
        """Ingest document chunks into Astra DB"""
        try:
            if not self.session:
                logger.error("Cannot ingest chunks: session is None")
                return False

            insert_stmt = self.session.prepare("""
                INSERT INTO document_chunks (chunk_id, content, embedding, metadata, document_type, filename, created_at)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            """)

            batch_size = 100
            successful_inserts = 0
            
            for i in range(0, len(chunks), batch_size):
                batch_chunks = chunks[i:i + batch_size]
                
                try:
                    for chunk in batch_chunks:
                        self.session.execute(insert_stmt, [
                            chunk.chunk_id,
                            chunk.content,
                            chunk.embedding,
                            json.dumps(chunk.metadata),
                            chunk.metadata.get('document_type', 'unknown'),
                            chunk.metadata.get('filename', 'unknown'),
                            datetime.now()
                        ])
                        successful_inserts += 1
                    
                    logger.info(f"Ingested batch {i//batch_size + 1}: {len(batch_chunks)} chunks")
                    
                except Exception as e:
                    logger.error(f"Error ingesting batch {i//batch_size + 1}: {e}")
                    continue

            logger.info(f"Successfully ingested {successful_inserts} out of {len(chunks)} chunks into Astra DB")
            return successful_inserts > 0
        
        except Exception as e:
            logger.error(f"Error ingesting chunks: {e}")
            return False

    def store_document_metadata(self, document_info: Dict[str, Any]) -> bool:
        """Store document metadata"""
        try:
            if not self.session:
                logger.error("Cannot store metadata: session is None")
                return False

            insert_stmt = self.session.prepare("""
                INSERT INTO document_metadata (document_id, filename, document_type, file_size, processed_at, total_chunks, processing_status)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            """)

            self.session.execute(insert_stmt, [
                document_info['document_id'],
                document_info['filename'],
                document_info['document_type'],
                document_info['file_size'],
                datetime.now(),
                document_info['total_chunks'],
                document_info['processing_status']
            ])

            return True
        except Exception as e:
            logger.error(f"Error storing document metadata: {e}")
            return False

    def get_chunk_count(self) -> int:
        """Get total number of chunks in the database"""
        try:
            if self.session:
                result = self.session.execute("SELECT COUNT(*) FROM document_chunks")
                return result.one()[0]
            else:
                logger.error("Cannot get chunk count: session is None")
                return 0
        except Exception as e:
            logger.error(f"Error getting chunk count: {e}")
            return 0

    def close(self):
        """Close database connection"""
        try:
            if self.session:
                self.session.shutdown()
            if self.cluster:
                self.cluster.shutdown()
            logger.info("Database connection closed")
        except Exception as e:
            logger.error(f"Error closing database connection: {e}")

def ingest_documents(file_paths: List[str], astra_config: Dict[str, str]) -> bool:
    """Main function to ingest multiple documents"""
    ingester = None
    try:
        # Initialize processor and ingester
        processor = DocumentProcessor()
        ingester = AstraDBIngester(
            token=astra_config['token'],
            database_id=astra_config['database_id'],
            region=astra_config.get('region', 'us-east1'),
            keyspace=astra_config.get('keyspace', 'insurance_docs')
        )

        all_chunks = []
        document_metadata = []

        # Process each document
        for file_path in file_paths:
            if not os.path.exists(file_path):
                logger.warning(f"File not found: {file_path}")
                continue

            logger.info(f"Processing document: {file_path}")
            chunks = processor.process_document(file_path)
            
            if chunks:
                all_chunks.extend(chunks)
                
                # Store document metadata
                doc_metadata = {
                    'document_id': str(uuid.uuid4()),
                    'filename': Path(file_path).name,
                    'document_type': chunks[0].metadata.get('document_type', 'unknown'),
                    'file_size': os.path.getsize(file_path),
                    'total_chunks': len(chunks),
                    'processing_status': 'completed'
                }
                document_metadata.append(doc_metadata)

        # Ingest all chunks
        if all_chunks:
            success = ingester.ingest_chunks(all_chunks)
            
            if success:
                # Store document metadata
                for doc_info in document_metadata:
                    ingester.store_document_metadata(doc_info)
                
                logger.info(f"Successfully ingested {len(all_chunks)} chunks from {len(file_paths)} documents")
                return True

        return False

    except Exception as e:
        logger.error(f"Error in document ingestion: {e}")
        return False
    finally:
        if ingester:
            ingester.close()

def setup_logging(log_level: str = "INFO"):
    """Setup logging configuration"""
    logging.basicConfig(
        level=getattr(logging, log_level.upper()),
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler('document_ingestion.log'),
            logging.StreamHandler()
        ]
    )

# Example usage
if __name__ == "__main__":
    # Setup logging
    setup_logging("INFO")

    # Example configuration
    astra_config = {
        'token': os.getenv('ASTRA_DB_APPLICATION_TOKEN'),
        'database_id': os.getenv('ASTRA_DB_ID'),
        'region': os.getenv('ASTRA_DB_REGION', 'us-east1'),
        'keyspace': os.getenv('ASTRA_DB_KEYSPACE', 'insurance_docs')
    }

    # Validate configuration
    if not astra_config['token'] or not astra_config['database_id']:
        logger.error("Missing required environment variables: ASTRA_DB_APPLICATION_TOKEN, ASTRA_DB_ID")
        exit(1)

    # Example file paths
    file_paths = [
        'data/documents/policy1.pdf',
        'data/documents/policy2.docx',
        'data/documents/terms.pptx'
    ]

    # Filter existing files
    existing_files = [path for path in file_paths if os.path.exists(path)]
    
    if not existing_files:
        logger.warning("No files found to process")
        exit(1)

    logger.info(f"Starting ingestion of {len(existing_files)} files")
    success = ingest_documents(existing_files, astra_config)
    
    if success:
        logger.info("Document ingestion completed successfully")
    else:
        logger.error("Document ingestion failed")
        exit(1)